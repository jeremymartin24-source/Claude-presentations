#!/usr/bin/env python3
"""
generate_deck.py — Interactive PowerPoint generator in the IT260 lecture style.

Prompts you step-by-step for module info, sections, and content, then shows
a structured JSON deck plan for you to confirm or edit before building the .pptx.

No API keys required.

Usage:
    python generate_deck.py
"""
import json
import sys
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Design constants ───────────────────────────────────────────────────────────
RED       = RGBColor(0x68, 0x00, 0x01)   # #680001  primary brand red
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
PINK      = RGBColor(0xF5, 0xC0, 0xC0)   # light pink — used on dark bg
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)   # body text on white bg
CARD_BG   = RGBColor(0xF5, 0xF5, 0xF5)   # overview / key-term card background
CARD_BDR  = RGBColor(0xDD, 0xDD, 0xDD)   # card border

SLIDE_W = Emu(9144000)
SLIDE_H = Emu(5143500)
FONT    = "Calibri"

HEADER_H  = Emu(1005840)   # height of the red banner on white-bg slides
MARGIN    = Emu(365760)    # ~0.4 in
INNER_MAR = Emu(411480)    # intro/body left margin
GAP       = Emu(91440)     # inter-element gap


# ── Pptx helpers ──────────────────────────────────────────────────────────────

def new_presentation():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def add_rect(slide, x, y, w, h, fill_rgb, border_rgb=None, border_w=Emu(12700)):
    from pptx.util import Emu as E
    shape = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    line = shape.line
    if border_rgb:
        line.color.rgb = border_rgb
        line.width = border_w
    else:
        line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text, size_pt, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, word_wrap=True, italic=False, space_before=0):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None


    p = tf.paragraphs[0]
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.name  = FONT
    run.font.size  = Pt(size_pt)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_footer(slide, course_label, university):
    """Dark red footer bar at the bottom of white-bg slides."""
    footer_h = Emu(228600)
    y = SLIDE_H - footer_h
    add_rect(slide, Emu(0), y, SLIDE_W, footer_h, RED)
    # left text
    add_textbox(slide, MARGIN, y, Emu(4000000), footer_h,
                course_label, 9, color=WHITE, align=PP_ALIGN.LEFT)
    # right text
    add_textbox(slide, Emu(5144000), y, Emu(3800000), footer_h,
                university, 9, color=WHITE, align=PP_ALIGN.RIGHT)


def add_red_header(slide, title_text, course_label, university):
    """Full-width red header banner + slide title + footer for content slides."""
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, HEADER_H, RED)
    add_textbox(slide, MARGIN, Emu(91440), Emu(7000000), HEADER_H - Emu(91440),
                title_text, 26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_footer(slide, course_label, university)


# ── Slide builders ─────────────────────────────────────────────────────────────

def build_title_slide(prs, data, course, professor):
    """
    Slide type: "title"
    data keys: title, subtitle
    """
    slide = blank_slide(prs)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RED

    # White header bar (top strip)
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, Emu(1051560), WHITE)

    # Course label (right-aligned, on white bar)
    add_textbox(slide, Emu(2926080), Emu(256032), Emu(5943600), Emu(502920),
                course, 13, color=DARK_GRAY, align=PP_ALIGN.RIGHT)

    # Module badge
    add_rect(slide, MARGIN, Emu(1389888), Emu(1554480), Emu(384048),
             RGBColor(0xFF, 0xFF, 0xFF),   # white, will look translucent-ish
             border_rgb=WHITE)
    add_textbox(slide, MARGIN, Emu(1389888), Emu(1554480), Emu(384048),
                data.get("badge", ""), 11, bold=True, color=RED, align=PP_ALIGN.CENTER)

    # Main title
    add_textbox(slide, MARGIN, Emu(1828800), Emu(8229600), Emu(1500000),
                data["title"], 48, bold=True, color=WHITE)

    # Subtitle
    add_textbox(slide, MARGIN, Emu(3350000), Emu(8229600), Emu(700000),
                data.get("subtitle", ""), 15, color=PINK)

    # Professor / institution line
    add_textbox(slide, MARGIN, Emu(4500000), Emu(8229600), Emu(400000),
                professor, 12, color=WHITE)

    return slide


def build_section_divider(prs, data, *_):
    """
    Slide type: "section_divider"
    data keys: label, title, subtitle
    """
    slide = blank_slide(prs)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RED

    # Small label (e.g. "Module 2")
    add_textbox(slide, MARGIN, Emu(1200000), Emu(7000000), Emu(400000),
                data.get("label", ""), 15, color=PINK)

    # Big title
    add_textbox(slide, MARGIN, Emu(1650000), Emu(8229600), Emu(1800000),
                data["title"], 44, bold=True, color=WHITE)

    # Tagline
    add_textbox(slide, MARGIN, Emu(3500000), Emu(8229600), Emu(600000),
                data.get("subtitle", ""), 16, color=PINK)

    return slide


def build_overview(prs, data, course, professor):
    """
    Slide type: "overview"
    data keys: heading, intro, sections:[{title, description}]
    """
    slide = blank_slide(prs)
    add_red_header(slide, data["heading"], course, professor)

    # Intro paragraph
    add_textbox(slide, INNER_MAR, HEADER_H + GAP,
                SLIDE_W - INNER_MAR * 2, Emu(500000),
                data.get("intro", ""), 13, color=DARK_GRAY)

    # Topic cards in a grid (2 columns)
    sections = data.get("sections", [])
    card_w = Emu(4206240)
    card_h = Emu(1200000)
    top_stripe_h = Emu(54864)
    col_gap = Emu(182880)
    row_gap = Emu(91440)
    start_y = HEADER_H + Emu(620000)
    start_x = MARGIN

    for i, sec in enumerate(sections[:6]):
        col = i % 2
        row = i // 2
        x = start_x + col * (card_w + col_gap)
        y = start_y + row * (card_h + row_gap)

        # Card background
        add_rect(slide, x, y, card_w, card_h, CARD_BG, CARD_BDR)
        # Red top stripe
        add_rect(slide, x, y, card_w, top_stripe_h, RED)
        # Section heading
        add_textbox(slide, x + GAP, y + top_stripe_h + GAP,
                    card_w - GAP * 2, Emu(200000),
                    sec["title"], 13, bold=True, color=RED)
        # Section description
        add_textbox(slide, x + GAP, y + top_stripe_h + Emu(240000),
                    card_w - GAP * 2, card_h - top_stripe_h - Emu(260000),
                    sec["description"], 11.5, color=DARK_GRAY)

    return slide


def build_objectives(prs, data, course, professor):
    """
    Slide type: "objectives"
    data keys: heading, intro, items:[str]
    """
    slide = blank_slide(prs)
    add_red_header(slide, data["heading"], course, professor)

    add_textbox(slide, INNER_MAR, HEADER_H + GAP,
                SLIDE_W - INNER_MAR * 2, Emu(380000),
                data.get("intro", "By the end of this module, you will be able to:"),
                13.5, color=DARK_GRAY)

    items = data.get("items", [])
    item_h = Emu(480000)
    badge_w = Emu(300000)
    badge_gap = Emu(60000)
    text_w = SLIDE_W - MARGIN * 2 - badge_w - badge_gap
    start_y = HEADER_H + Emu(480000)

    # Two-column layout for items
    col_count = 2 if len(items) > 3 else 1
    col_w = (SLIDE_W - MARGIN * 2 - GAP * (col_count - 1)) // col_count

    for i, item in enumerate(items[:6]):
        col = i % col_count if col_count == 2 else 0
        row = i // col_count if col_count == 2 else i
        x = MARGIN + col * (col_w + GAP)
        y = start_y + row * (item_h + GAP)

        # Number badge (red square)
        add_rect(slide, x, y, badge_w, badge_w, RED)
        num_str = f"{(i + 1):02d}"
        add_textbox(slide, x, y, badge_w, badge_w,
                    num_str, 13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Item text
        add_textbox(slide, x + badge_w + badge_gap, y,
                    col_w - badge_w - badge_gap, item_h,
                    item, 13, color=DARK_GRAY)

    return slide


def build_content_bullets(prs, data, course, professor):
    """
    Slide type: "content_bullets"
    Single column: intro paragraph + sections with red headings and bullet lists.
    data keys: heading, intro, sections:[{heading, bullets:[str]}]
    """
    slide = blank_slide(prs)
    add_red_header(slide, data["heading"], course, professor)

    y = HEADER_H + GAP
    if data.get("intro"):
        add_textbox(slide, INNER_MAR, y,
                    SLIDE_W - INNER_MAR * 2, Emu(480000),
                    data["intro"], 12.5, color=DARK_GRAY)
        y += Emu(500000)

    for sec in data.get("sections", []):
        # Red heading
        add_textbox(slide, INNER_MAR, y,
                    SLIDE_W - INNER_MAR * 2, Emu(240000),
                    sec["heading"], 14, bold=True, color=RED)
        y += Emu(250000)
        # Bullets
        for bullet in sec.get("bullets", []):
            add_textbox(slide, INNER_MAR + Emu(60000), y,
                        SLIDE_W - INNER_MAR * 2 - Emu(60000), Emu(200000),
                        f"▸  {bullet}", 10.5, color=DARK_GRAY)
            y += Emu(210000)
        y += GAP

    return slide


def build_content_columns(prs, data, course, professor):
    """
    Slide type: "content_columns"
    2-3 column layout, each column has a red header + bullet list.
    data keys: heading, intro, columns:[{heading, bullets:[str]}]
    """
    slide = blank_slide(prs)
    add_red_header(slide, data["heading"], course, professor)

    y = HEADER_H + GAP
    if data.get("intro"):
        add_textbox(slide, INNER_MAR, y,
                    SLIDE_W - INNER_MAR * 2, Emu(420000),
                    data["intro"], 12.5, color=DARK_GRAY)
        y += Emu(440000)

    cols = data.get("columns", [])
    n = max(1, min(len(cols), 3))
    avail_w = SLIDE_W - MARGIN * 2
    col_w = (avail_w - GAP * (n - 1)) // n
    avail_h = SLIDE_H - y - Emu(300000)

    col_header_h = Emu(240000)

    for i, col in enumerate(cols[:3]):
        x = MARGIN + i * (col_w + GAP)
        # Column header bar
        add_rect(slide, x, y, col_w, col_header_h, RED)
        add_textbox(slide, x + GAP, y, col_w - GAP, col_header_h,
                    col["heading"], 14, bold=True, color=WHITE)
        # Bullets
        by = y + col_header_h + GAP
        for bullet in col.get("bullets", []):
            bh = Emu(200000)
            add_textbox(slide, x + Emu(30000), by,
                        col_w - Emu(30000), bh,
                        f"▸  {bullet}", 10.5, color=DARK_GRAY)
            by += bh + Emu(30000)

    return slide


def build_content_table(prs, data, course, professor):
    """
    Slide type: "content_table"
    A scenario/comparison table.
    data keys: heading, intro, rows:[{label, normal, threat, impact}]
    OR generic: columns:[str], rows:[[str]]
    """
    from pptx.util import Pt as P
    slide = blank_slide(prs)
    add_red_header(slide, data["heading"], course, professor)

    y = HEADER_H + GAP
    if data.get("intro"):
        add_textbox(slide, INNER_MAR, y,
                    SLIDE_W - INNER_MAR * 2, Emu(380000),
                    data["intro"], 12.5, color=DARK_GRAY)
        y += Emu(400000)

    col_headers = data.get("columns", [])
    rows = data.get("rows", [])
    if not col_headers or not rows:
        return slide

    n_cols = len(col_headers)
    avail_w = SLIDE_W - MARGIN * 2
    col_w = avail_w // n_cols
    row_h = Emu(350000)
    hdr_h = Emu(280000)

    # Header row
    for ci, hdr in enumerate(col_headers):
        x = MARGIN + ci * col_w
        add_rect(slide, x, y, col_w, hdr_h, RED, WHITE, Emu(6350))
        add_textbox(slide, x + GAP, y, col_w - GAP, hdr_h,
                    hdr, 12, bold=True, color=WHITE)
    y += hdr_h

    # Data rows
    for ri, row in enumerate(rows):
        bg = CARD_BG if ri % 2 == 0 else WHITE
        for ci, cell in enumerate(row[:n_cols]):
            x = MARGIN + ci * col_w
            add_rect(slide, x, y, col_w, row_h, bg, CARD_BDR, Emu(6350))
            add_textbox(slide, x + GAP, y + GAP, col_w - GAP * 2,
                        row_h - GAP, cell, 10.5, color=DARK_GRAY)
        y += row_h

    return slide


def build_key_terms(prs, data, course, professor):
    """
    Slide type: "key_terms"
    data keys: heading, terms:[{term, definition}]
    """
    slide = blank_slide(prs)
    add_red_header(slide, data["heading"], course, professor)

    terms = data.get("terms", [])
    col_count = 2 if len(terms) > 4 else 1
    avail_w = SLIDE_W - MARGIN * 2
    col_w = (avail_w - (GAP if col_count == 2 else 0)) // col_count
    item_h = Emu(380000)
    start_y = HEADER_H + Emu(100000)

    for i, t in enumerate(terms[:8]):
        col = i % col_count
        row = i // col_count
        x = MARGIN + col * (col_w + GAP)
        y = start_y + row * (item_h + GAP // 2)

        # Term (red, bold)
        add_textbox(slide, x, y, col_w, Emu(200000),
                    t["term"], 12, bold=True, color=RED)
        # Definition (gray)
        add_textbox(slide, x, y + Emu(200000), col_w, Emu(200000),
                    t["definition"], 11.5, color=DARK_GRAY)

    return slide


def build_case_study(prs, data, course, professor):
    """
    Slide type: "case_study"
    A real-world example with labelled panels.
    data keys: heading, intro, panels:[{label, content}]
    """
    slide = blank_slide(prs)
    add_red_header(slide, data["heading"], course, professor)

    y = HEADER_H + GAP
    if data.get("intro"):
        add_textbox(slide, INNER_MAR, y,
                    SLIDE_W - INNER_MAR * 2, Emu(380000),
                    data["intro"], 12.5, color=DARK_GRAY)
        y += Emu(400000)

    panels = data.get("panels", [])
    label_w = Emu(1400000)
    content_w = SLIDE_W - MARGIN * 2 - label_w - GAP
    panel_h = Emu(350000)

    for panel in panels:
        # Label box (red)
        add_rect(slide, MARGIN, y, label_w, panel_h, RED)
        add_textbox(slide, MARGIN + Emu(30000), y, label_w - Emu(60000), panel_h,
                    panel["label"], 12, bold=True, color=WHITE)
        # Content box
        add_rect(slide, MARGIN + label_w + GAP, y, content_w, panel_h, CARD_BG, CARD_BDR)
        add_textbox(slide, MARGIN + label_w + GAP + Emu(30000), y + Emu(30000),
                    content_w - Emu(60000), panel_h - Emu(30000),
                    panel["content"], 11, color=DARK_GRAY)
        y += panel_h + GAP // 2

    return slide


def build_takeaways(prs, data, course, professor):
    """
    Slide type: "takeaways"
    data keys: heading, items:[str], next_module (optional)
    """
    slide = blank_slide(prs)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RED

    # Heading
    add_textbox(slide, MARGIN, Emu(150000), SLIDE_W - MARGIN * 2, Emu(400000),
                data["heading"], 28, bold=True, color=WHITE)

    items = data.get("items", [])
    badge_w = Emu(270000)
    item_h  = Emu(330000)
    text_w  = SLIDE_W - MARGIN * 2 - badge_w - GAP
    start_y = Emu(600000)

    col_count = 2 if len(items) > 4 else 1
    if col_count == 2:
        col_w = (SLIDE_W - MARGIN * 2 - GAP) // 2

    for i, item in enumerate(items[:8]):
        col = i % col_count if col_count == 2 else 0
        row = i // col_count if col_count == 2 else i
        x = MARGIN + col * (col_w + GAP) if col_count == 2 else MARGIN
        w = col_w if col_count == 2 else text_w + badge_w + GAP
        y = start_y + row * (item_h + GAP // 2)

        # Number badge (white)
        add_rect(slide, x, y, badge_w, badge_w, WHITE)
        add_textbox(slide, x, y, badge_w, badge_w,
                    str(i + 1), 12, bold=True, color=RED, align=PP_ALIGN.CENTER)
        # Item text
        tw = (col_w if col_count == 2 else SLIDE_W - MARGIN * 2) - badge_w - GAP
        add_textbox(slide, x + badge_w + GAP, y, tw, item_h,
                    item, 12.5, color=WHITE)

    if data.get("next_module"):
        add_textbox(slide, MARGIN, SLIDE_H - Emu(350000),
                    SLIDE_W - MARGIN * 2, Emu(250000),
                    data["next_module"], 12, color=PINK, italic=True)

    return slide


def build_assignment(prs, data, course, professor):
    """
    Slide type: "assignment"
    data keys: heading, intro, items:[{number, text}]
    """
    slide = blank_slide(prs)
    add_red_header(slide, data["heading"], course, professor)

    y = HEADER_H + GAP
    if data.get("intro"):
        add_textbox(slide, INNER_MAR, y,
                    SLIDE_W - INNER_MAR * 2, Emu(400000),
                    data["intro"], 13, color=DARK_GRAY)
        y += Emu(430000)

    for item in data.get("items", []):
        # Number badge
        badge_w = Emu(270000)
        add_rect(slide, INNER_MAR, y, badge_w, badge_w, RED)
        add_textbox(slide, INNER_MAR, y, badge_w, badge_w,
                    str(item["number"]), 13, bold=True,
                    color=WHITE, align=PP_ALIGN.CENTER)
        # Text
        add_textbox(slide, INNER_MAR + badge_w + GAP, y,
                    SLIDE_W - INNER_MAR * 2 - badge_w - GAP, Emu(280000),
                    item["text"], 13, color=DARK_GRAY)
        y += Emu(290000) + GAP // 2

    return slide


# ── Slide dispatch ─────────────────────────────────────────────────────────────

BUILDERS = {
    "title":            build_title_slide,
    "section_divider":  build_section_divider,
    "overview":         build_overview,
    "objectives":       build_objectives,
    "content_bullets":  build_content_bullets,
    "content_columns":  build_content_columns,
    "content_table":    build_content_table,
    "key_terms":        build_key_terms,
    "case_study":       build_case_study,
    "takeaways":        build_takeaways,
    "assignment":       build_assignment,
}


def build_slide(prs, slide_data, course, professor):
    stype = slide_data.get("type")
    builder = BUILDERS.get(stype)
    if builder is None:
        print(f"  ⚠  Unknown slide type '{stype}' — skipping.", file=sys.stderr)
        return
    builder(prs, slide_data, course, professor)



# ── Interactive prompt helpers ─────────────────────────────────────────────────

def ask(prompt_text, default=None):
    """Single-line prompt. Returns stripped input or default."""
    if default:
        display = f"{prompt_text} [{default}]: "
    else:
        display = f"{prompt_text}: "
    val = input(display).strip()
    return val if val else (default or "")


def ask_list(prompt_text, max_items=8, min_items=0):
    """Collect a list of lines until blank. Returns list of non-empty strings."""
    print(f"\n{prompt_text}")
    print("  (press Enter on a blank line to finish)")
    items = []
    while len(items) < max_items:
        line = input(f"  {len(items)+1}. ").strip()
        if not line:
            if len(items) >= min_items:
                break
            print(f"  Need at least {min_items} item(s).")
        else:
            items.append(line)
    return items


def section_break(title):
    print(f"\n{'─'*50}")
    print(f"  {title}")
    print(f"{'─'*50}")


# ── Slide-type heuristics ──────────────────────────────────────────────────────

COLUMN_WORDS  = {"types", "categories", "kinds", "forms", "pillars",
                 "components", "elements", "three", "two", "four", "five"}
TABLE_WORDS   = {"vs", "versus", "compare", "comparison", "difference",
                 "contrast", "between", "against"}
CASE_WORDS    = {"example", "case", "scenario", "incident", "attack",
                 "breach", "failure", "real", "event", "story"}

def detect_slide_type(title):
    words = set(title.lower().split())
    if words & TABLE_WORDS:
        return "content_table"
    if words & COLUMN_WORDS:
        return "content_columns"
    if words & CASE_WORDS:
        return "case_study"
    return "content_bullets"


TYPE_MENU = """\
  Content type for this section:
    1. Narrative explanation  (bullet lists under sub-headings)
    2. Side-by-side columns   (2–3 parallel columns)
    3. Structured table       (rows and columns)
    4. Real-world case study  (labelled panels)
    5. Auto-detect from title
"""

TYPE_MAP = {
    "1": "content_bullets",
    "2": "content_columns",
    "3": "content_table",
    "4": "case_study",
    "5": "auto",
}


# ── Section content gatherers ──────────────────────────────────────────────────

def gather_bullets(section_title):
    """Return content_bullets slide data."""
    print(f"\n  Gathering bullet content for: {section_title}")
    num_subs = int(ask("  How many sub-headings?", "2"))
    sections = []
    for i in range(max(1, min(num_subs, 4))):
        heading = ask(f"  Sub-heading {i+1}")
        bullets = ask_list(f"  Bullets under '{heading}'", max_items=6, min_items=1)
        sections.append({"heading": heading, "bullets": bullets})
    intro = ask("  Optional intro sentence (or Enter to skip)")
    return {
        "type": "content_bullets",
        "heading": section_title,
        "intro": intro,
        "sections": sections,
    }


def gather_columns(section_title):
    """Return content_columns slide data."""
    print(f"\n  Gathering column content for: {section_title}")
    num_cols = int(ask("  How many columns? (2 or 3)", "2"))
    columns = []
    for i in range(max(2, min(num_cols, 3))):
        heading = ask(f"  Column {i+1} heading")
        bullets = ask_list(f"  Bullets under '{heading}'", max_items=6, min_items=1)
        columns.append({"heading": heading, "bullets": bullets})
    intro = ask("  Optional intro sentence (or Enter to skip)")
    return {
        "type": "content_columns",
        "heading": section_title,
        "intro": intro,
        "columns": columns,
    }


def gather_table(section_title):
    """Return content_table slide data."""
    print(f"\n  Gathering table content for: {section_title}")
    col_input = ask("  Column headers (comma-separated)", "Term, Definition, Example")
    columns = [c.strip() for c in col_input.split(",") if c.strip()]
    intro = ask("  Optional intro sentence (or Enter to skip)")
    print(f"\n  Enter table rows. Each row: comma-separated values ({len(columns)} columns).")
    print("  (blank line to finish)")
    rows = []
    while True:
        row_input = input(f"  Row {len(rows)+1}: ").strip()
        if not row_input:
            if rows:
                break
            print("  Need at least one row.")
            continue
        cells = [c.strip() for c in row_input.split(",")]
        # Pad or trim to match column count
        while len(cells) < len(columns):
            cells.append("")
        rows.append(cells[:len(columns)])
    return {
        "type": "content_table",
        "heading": section_title,
        "intro": intro,
        "columns": columns,
        "rows": rows,
    }


def gather_case_study(section_title):
    """Return case_study slide data."""
    print(f"\n  Gathering case study content for: {section_title}")
    intro = ask("  One-sentence framing (what this example illustrates)")
    print("\n  Enter panels (label + content). Blank label to finish. Suggested labels:")
    print("    The Threat / The Vulnerability / ⚠ The Impact  — or name them yourself.")
    panels = []
    while len(panels) < 5:
        label = input(f"\n  Panel {len(panels)+1} label (or blank to finish): ").strip()
        if not label and panels:
            break
        elif not label:
            print("  Need at least one panel.")
            continue
        content = ask(f"  Panel '{label}' content")
        panels.append({"label": label, "content": content})
    return {
        "type": "case_study",
        "heading": section_title,
        "intro": intro,
        "panels": panels,
    }


GATHERERS = {
    "content_bullets": gather_bullets,
    "content_columns": gather_columns,
    "content_table":   gather_table,
    "case_study":      gather_case_study,
}


# ── Plan builder ───────────────────────────────────────────────────────────────

def build_plan(meta, objectives, sections, key_terms, assignment, takeaways, next_mod):
    """Assemble the full slides list from gathered data."""
    mod = meta["module"]
    slides = []

    # 1. Title
    slides.append({
        "type": "title",
        "badge": f"MODULE {mod}",
        "title": meta["title"],
        "subtitle": meta["subtitle"],
    })

    # 2. Overview — built from section titles + descriptions
    overview_sections = [
        {"title": s["title"], "description": s["description"]}
        for s in sections
    ]
    slides.append({
        "type": "overview",
        "heading": f"Module {mod} · Overview",
        "intro": meta.get("overview_intro", ""),
        "sections": overview_sections,
    })

    # 3. Objectives
    slides.append({
        "type": "objectives",
        "heading": f"Module {mod} · Learning Objectives",
        "intro": "By the end of this module, you will be able to:",
        "items": objectives,
    })

    # 4. Sections — each gets a divider + content slide
    for sec in sections:
        slides.append({
            "type": "section_divider",
            "label": f"Module {mod}",
            "title": sec["title"],
            "subtitle": sec["description"],
        })
        slides.append(sec["content"])

    # 5. Key terms
    if key_terms:
        slides.append({
            "type": "key_terms",
            "heading": f"Module {mod} · Key Terms",
            "terms": key_terms,
        })

    # 6. Assignment
    if assignment:
        slides.append(assignment)

    # 7. Takeaways
    slides.append({
        "type": "takeaways",
        "heading": f"Module {mod} · Key Takeaways",
        "items": takeaways,
        "next_module": f"Next: {next_mod}" if next_mod else "",
    })

    return {"meta": meta, "slides": slides}


# ── Confirmation / edit loop ───────────────────────────────────────────────────

def confirm_plan(plan):
    """Show plan as JSON, let user confirm or paste edits. Returns final plan dict."""
    json_str = json.dumps(plan, indent=2)
    print("\n" + "═"*60)
    print("  DECK PLAN")
    print("═"*60)
    print(json_str)
    print("═"*60)
    print("\nConfirm this plan?")
    print("  Y / Enter  — build the deck")
    print("  q          — quit")
    print("  (paste)    — paste edited JSON, then type END on its own line\n")

    while True:
        response = input("> ").strip()
        if response.lower() in ("y", ""):
            return plan
        if response.lower() == "q":
            print("Aborted.")
            sys.exit(0)
        # Treat as start of pasted JSON
        lines = [response]
        print("Paste your edited JSON. Type END on its own line when done.")
        while True:
            line = input()
            if line.strip() == "END":
                break
            lines.append(line)
        raw = "\n".join(lines)
        try:
            edited = json.loads(raw)
            print("✓ Parsed edited JSON.")
            return edited
        except json.JSONDecodeError as e:
            print(f"JSON parse error: {e}")
            print("Try again — Y to use original, q to quit, or paste again.")


# ── Main interactive flow ──────────────────────────────────────────────────────

def main():
    print("\n" + "═"*60)
    print("  IT260 Lecture Deck Generator")
    print("═"*60)

    # ── Step 1: Basic info ─────────────────────────────────────────────────
    section_break("Step 1 of 6 — Basic Info")
    course    = ask("Course name", "IT260 · IT Risk Management")
    mod_num   = ask("Module number")
    title     = ask("Module title")
    subtitle  = ask("Module tagline (one sentence)")
    professor = ask("Professor · University · Term",
                    "Professor Martin · University of Northwestern Ohio · Fall 2025")
    output    = ask("Output filename", f"Module{mod_num}.pptx")
    ov_intro  = ask("Overview slide intro sentence (summarises the module in 1–2 lines)")

    meta = {
        "course": course,
        "module": mod_num,
        "title": title,
        "subtitle": subtitle,
        "professor": professor,
        "output": output,
        "overview_intro": ov_intro,
    }

    # ── Step 2: Learning objectives ────────────────────────────────────────
    section_break("Step 2 of 6 — Learning Objectives")
    print("  Action-verb sentences (Describe…, Define…, Explain…, Apply…)")
    objectives = ask_list("Enter objectives", max_items=6, min_items=2)

    # ── Step 3: Sections ───────────────────────────────────────────────────
    section_break("Step 3 of 6 — Topic Sections")
    num_sections = int(ask("How many major topic sections?", "3"))
    sections = []
    for i in range(max(1, min(num_sections, 5))):
        print(f"\n  ── Section {i+1} of {num_sections} ──")
        sec_title = ask("  Section title")
        sec_desc  = ask("  One-sentence description (shown on Overview slide)")

        print(TYPE_MENU)
        choice = ask("  Choice", "5").strip()
        stype = TYPE_MAP.get(choice, "auto")
        if stype == "auto":
            stype = detect_slide_type(sec_title)
            print(f"  → Auto-detected type: {stype}")

        gatherer = GATHERERS.get(stype, gather_bullets)
        content_slide = gatherer(sec_title)

        sections.append({
            "title": sec_title,
            "description": sec_desc,
            "content": content_slide,
        })

    # ── Step 4: Key terms ──────────────────────────────────────────────────
    section_break("Step 4 of 6 — Key Terms")
    key_terms = []
    if ask("Include a Key Terms slide? (y/n)", "y").lower() == "y":
        print("  Enter terms. Blank term to finish.")
        while len(key_terms) < 8:
            term = input(f"\n  Term {len(key_terms)+1} (or blank to finish): ").strip()
            if not term:
                break
            definition = ask(f"  Definition for '{term}'")
            key_terms.append({"term": term, "definition": definition})

    # ── Step 5: Assignment ─────────────────────────────────────────────────
    section_break("Step 5 of 6 — Assignment")
    assignment = None
    if ask("Include an Assignment slide? (y/n)", "y").lower() == "y":
        a_intro = ask("Assignment intro sentence")
        a_items = []
        print("  Enter assignment prompts. Blank to finish.")
        while len(a_items) < 4:
            text = input(f"\n  Prompt {len(a_items)+1} (or blank to finish): ").strip()
            if not text:
                if a_items:
                    break
            else:
                a_items.append({"number": len(a_items)+1, "text": text})
        assignment = {
            "type": "assignment",
            "heading": "This Week — Focus and Assignment",
            "intro": a_intro,
            "items": a_items,
        }

    # ── Step 6: Takeaways ──────────────────────────────────────────────────
    section_break("Step 6 of 6 — Key Takeaways")
    print("  Enter 4–8 takeaway statements (one per line).")
    takeaways = ask_list("Takeaways", max_items=8, min_items=4)
    next_mod  = ask("Next module title (or blank to skip)")

    # ── Build plan & confirm ───────────────────────────────────────────────
    plan = build_plan(meta, objectives, sections, key_terms, assignment, takeaways, next_mod)
    plan = confirm_plan(plan)

    # ── Render PPTX ───────────────────────────────────────────────────────
    slides   = plan.get("slides", [])
    out_meta = plan.get("meta", meta)
    course_l = out_meta.get("course", course)
    prof_l   = out_meta.get("professor", professor)
    out_file = out_meta.get("output", output)

    print(f"\nBuilding {len(slides)} slides…")
    prs = new_presentation()
    for i, slide_data in enumerate(slides):
        stype = slide_data.get("type", "?")
        print(f"  [{i+1}/{len(slides)}] {stype}")
        build_slide(prs, slide_data, course_l, prof_l)

    prs.save(out_file)
    print(f"\n✓ Saved: {out_file}")


if __name__ == "__main__":
    main()
